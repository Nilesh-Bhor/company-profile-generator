import re
import os
import io
import json
import pdfkit
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.utility import get_logo
from agents.agent_factory import AgentFactory, AgentType

class CompanyProfile:
    def __init__(self, company_name, company_website=None, agent_type=None):
        self.profile_data = None
        self.markdown_data = None
        self.company_name = company_name
        self.company_website = company_website
        logo_name = company_name.lower().replace(' ', '').replace('.', '')
        self.logo_url = f"https://logo.clearbit.com/{logo_name}.com"
        self.logo = get_logo(self.logo_url)

        if agent_type is None:
            # Get agent type from environment variable
            agent_type_str = os.getenv('AGENT_TYPE', 'GOOGLE_GEMINI')
            agent_type = AgentType[agent_type_str] if agent_type_str in AgentType.__members__ else AgentType.GOOGLE_GEMINI
        
        # Configure the agent
        self.agent = AgentFactory.get_agent(agent_type)
    

    def create_markdown(self):
        markdown_text = ""
        try:
            if self.profile_data is not None:
                # Company Overview
                if 'company_overview' in self.profile_data:
                    overview = self.profile_data['company_overview']
            
                company_name = overview.get('name', self.company_name)
                logo_url = self.logo_url if self.logo is not None else overview['logo_url']

                self.profile_data['company_overview']['logo_url'] = logo_url
                
                markdown_text += f"# <img src='{logo_url}' height='50' align='left' style='margin-right: 10px;'> {company_name} \n\n"
                markdown_text += f"{overview.get('description', '')}\n\n"
                
                for key in ['industry', 'location', 'mission', 'vision']:
                    if key in overview:
                        markdown_text += f"* **{key.title()}**: {overview[key]}\n"
                
                markdown_text += "\n---\n\n"
        
                # Products and Services
                if 'products_and_services' in self.profile_data:
                    # Products
                    if 'products' in self.profile_data['products_and_services']:
                        markdown_text += "## Products\n\n"
                        products = self.profile_data['products_and_services']['products']
                        markdown_text += f"{products.get('description', '')}\n\n"
                        for item in products.get('items', []):
                            markdown_text += f"#### {item['name']}\n"
                            markdown_text += f"{item['description']}\n\n"
                            if 'features' in item:
                                markdown_text += "**Key Features:**\n"
                                for feature in item['features']:
                                    markdown_text += f"* {feature}\n"
                            markdown_text += "\n"
                        
                        markdown_text += "---\n\n"

                    # Services
                    if 'services' in self.profile_data['products_and_services']:
                        markdown_text += "## Services\n\n"
                        services = self.profile_data['products_and_services']['services']
                        markdown_text += f"{services.get('description', '')}\n\n"
                        for item in services.get('items', []):
                            markdown_text += f"#### {item['name']}\n"
                            markdown_text += f"{item['description']}\n\n"
                            if 'features' in item:
                                markdown_text += "**Key Features:**\n"
                                for feature in item['features']:
                                    markdown_text += f"* {feature}\n"
                            markdown_text += "\n"
                        
                        markdown_text += "---\n\n"

                # Management Team
                if 'management_team' in self.profile_data:
                    markdown_text += "## Management Team\n\n"
                    team = self.profile_data['management_team']
                    markdown_text += f"{team.get('description', '')}\n\n"
                    for member in team.get('members', []):
                        markdown_text += f"#### {member['name']} - {member['position']}\n"
                        markdown_text += f"{member['qualifications']}\n\n"
                    
                    markdown_text += "---\n\n"
                
                # Milestones
                if 'milestones' in self.profile_data and self.profile_data['milestones']:
                    markdown_text += "## Key Milestones\n\n"
                    for milestone in self.profile_data['milestones']:
                        markdown_text += f"* **{milestone['date']}**: {milestone['description']}\n"
                    
                    markdown_text += "\n---\n\n"
                
                # Financial Highlights
                if 'financial_highlights' in self.profile_data:
                    markdown_text += "## Financial Highlights\n\n"
                    financials = self.profile_data['financial_highlights']
                    markdown_text += f"{financials.get('overview', '')}\n\n"
                    if 'metrics' in financials:
                        for metric in financials['metrics']:
                            markdown_text += f"* **{metric['year']}**: Revenue {metric['revenue']}, Growth {metric['growth']}\n"
                    
                    markdown_text += "\n---\n"
                
                # Sources
                if 'sources' in self.profile_data:
                    markdown_text += "## Sources\n\n"
                    for source in self.profile_data['sources']:
                        markdown_text += f"* {source}\n"
                    
                    markdown_text += "\n---\n"

        except Exception as e:
            print(f"An error occurred while generating markdown: {str(e)}")
        
        return markdown_text
    

    def get_company_profile(self):
        try:
            prompt = self.agent.get_company_profile_prompt(self.company_name, self.company_website)
            additional_info = self.agent.search_company_info(self.company_name)
            prompt += additional_info

            response_text = self.agent.generate_content(prompt)

            if len(response_text) > 0:
                # Clean up the response by removing markdown code block markers
                json_response = response_text.strip()
                json_response = re.sub(r'^.*?```json\s*|\s*```$', '', json_response, flags=re.DOTALL)

                self.profile_data = json.loads(json_response)
                self.markdown_data = self.create_markdown()

            return self.markdown_data
        except Exception as e:
            print(f"An error occurred while generating company profile: {str(e)}")
        
        return None
    
    
    def generate_ppt(self):
        if self.profile_data is not None:
            try:
                prs = Presentation()
                    
                # Title slide with logo
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                title = title_slide.shapes.title
                company_name = self.profile_data['company_overview']['name'] if 'company_overview' in self.profile_data else self.company_name
                title.text = company_name

                subtitle = title_slide.placeholders[1]
                company_site = self.profile_data['company_overview']['website'] if 'company_overview' in self.profile_data else self.company_website
                subtitle.text = f"{company_site} \n\n Company Profile" if company_site else "Company Profile"
                
                logo = self.logo
                if logo is None:
                    if 'company_overview' in self.profile_data and 'logo_url' in self.profile_data['company_overview']:
                        try:
                            logo_url = self.profile_data['company_overview']['logo_url']
                            logo = get_logo(logo_url)
                        except Exception as e:
                            print(f"Error adding logo: {str(e)}")
                
                if logo is not None:
                    left = Inches(4)  # Center position
                    top = Inches(0.5)  # Adjusted top position
                    width = Inches(2)  # Fixed width
                    height = Inches(2)  # Fixed height
                    title_slide.shapes.add_picture(logo, left, top, width, height)

                # Company Overview slide
                if 'company_overview' in self.profile_data:
                    overview = self.profile_data['company_overview']

                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = overview['name']
                    content = slide.placeholders[1].text_frame            
                    p = content.paragraphs[0]
                    p.text = overview['description']
                    p.font.size = Pt(14)
                    
                    for key in ['industry', 'location', 'mission', 'vision']:
                        if key in overview:
                            p = content.add_paragraph()
                            p.text = f"{key.title()}: {overview[key]}"
                            p.font.size = Pt(14)
                            p.level = 0
                
                # Products and Services
                if 'products_and_services' in self.profile_data:
                    # Products slide
                    if 'products' in self.profile_data['products_and_services']:
                        products = self.profile_data['products_and_services']['products']
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        title = slide.shapes.title
                        title.text = "Products"
                        
                        content = slide.placeholders[1].text_frame
                        p = content.paragraphs[0]
                        p.text = products['description']
                        p.font.size = Pt(14)

                        for i, item in enumerate(products['items']):
                            if i % 2 == 0 and i > 0:
                                slide = prs.slides.add_slide(prs.slide_layouts[1])
                                title = slide.shapes.title
                                title.text = "Products"
                                content = slide.placeholders[1].text_frame
                                p = content.paragraphs[0]
                            else:
                                p = content.add_paragraph()

                            p.text = f"{item['name']}: {item['description']}"
                            p.font.size = Pt(14)
                            p.level = 0

                            if 'features' in item:
                                p = content.add_paragraph()
                                p.text = "Key Features:"
                                p.font.size = Pt(14)
                                p.level = 1
                                for feature in item['features']:
                                    p = content.add_paragraph() 
                                    p.text = feature
                                    p.font.size = Pt(14)
                                    p.level = 2

                    # Services slide
                    if 'services' in self.profile_data['products_and_services']:
                        services = self.profile_data['products_and_services']['services']
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        title = slide.shapes.title
                        title.text = "Services"
                        
                        content = slide.placeholders[1].text_frame  
                        p = content.paragraphs[0]
                        p.text = services['description']
                        p.font.size = Pt(14)

                        for i, item in enumerate(services['items']):
                            if i % 2 == 0 and i > 0:
                                slide = prs.slides.add_slide(prs.slide_layouts[1])
                                title = slide.shapes.title
                                title.text = "Services"
                                content = slide.placeholders[1].text_frame
                                p = content.paragraphs[0]
                            else:
                                p = content.add_paragraph()

                            p.text = f"{item['name']}: {item['description']}"
                            p.font.size = Pt(14)
                            p.level = 0

                            if 'features' in item:
                                p = content.add_paragraph() 
                                p.text = "Key Features:"
                                p.font.size = Pt(14)
                                p.level = 1
                                for feature in item['features']:
                                    p = content.add_paragraph()
                                    p.text = feature
                                    p.font.size = Pt(14)
                                    p.level = 2
                
                # Management Team slide
                if 'management_team' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Management Team"
                    
                    content = slide.placeholders[1].text_frame
                    team = self.profile_data['management_team']
                    
                    p = content.paragraphs[0]
                    p.text = team['description']
                    p.font.size = Pt(14)
                    
                    for member in team['members']:
                        p = content.add_paragraph()
                        p.text = f"{member['name']} - {member['position']}"
                        p.font.size = Pt(14)
                        p.level = 0
                        
                        p = content.add_paragraph()
                        p.text = member['qualifications']
                        p.font.size = Pt(14)
                        p.level = 1
                
                # Milestones slide
                if 'milestones' in self.profile_data and self.profile_data['milestones']:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Key Milestones"

                    content = slide.placeholders[1].text_frame
                    content.paragraphs[0].font.size = Pt(14)
                    
                    for milestone in self.profile_data['milestones']:
                        p = content.add_paragraph()
                        p.text = f"{milestone['date']}: {milestone['description']}"
                        p.font.size = Pt(14)
                        p.level = 0
                
                # Financial Highlights slide
                if 'financial_highlights' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Financial Highlights"
                    
                    content = slide.placeholders[1].text_frame
                    financials = self.profile_data['financial_highlights']
                    
                    p = content.paragraphs[0]
                    p.text = financials['overview']
                    p.font.size = Pt(14)
                    
                    if 'metrics' in financials:
                        for metric in financials['metrics']:
                            p = content.add_paragraph()
                            p.text = f"{metric['year']}: Revenue {metric['revenue']}, Growth {metric['growth']}"
                            p.font.size = Pt(14)
                            p.level = 1                
                
                # Save to bytes buffer
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                return ppt_buffer
            except Exception as e:
                print(f"An error occurred while generating PPT: {str(e)}")
        
        return None
    

    def generate_pdf(self):
        try:
            if self.markdown_data is not None:
                wkhtmltopdf_path = os.getenv('WKHTMLTOPDF_PATH')
                config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)

                # Convert markdown to HTML
                html = markdown.markdown(self.markdown_data)
                
                # Parse with BeautifulSoup to ensure proper HTML structure
                soup = BeautifulSoup(html, 'html.parser')
                
                # Add basic styling
                style = soup.new_tag('style')
                style.string = """
                    body { font-family: Arial, sans-serif; line-height: 1.6; padding: 20px; }
                    h1, h2 { color: #333; }
                    hr { border: 1px solid #ddd; }
                    ul { margin-left: 20px; }
                """
                soup.insert(0, style)
                
                # Convert back to string
                html_content = str(soup)
                pdf_buffer = pdfkit.from_string(html_content, configuration=config)
                return pdf_buffer
        except Exception as e:
            print(f"An error occurred while generating PDF: {str(e)}")

        return None
    

    def format_profile(self, profile_data):
        """Format an existing profile from data"""
        self.profile_data = profile_data
        self.markdown_data = self.create_markdown()

        return self.markdown_data
    

    def update_profile_with_markdown(self, markdown_content):
        """Update the profile data based on edited markdown content"""
        # Store the original markdown
        self.profile = markdown_content
        
        # Here you could add logic to parse the markdown and update specific sections
        # of the profile_data if needed. For now, we'll just keep the markdown updated
        # and regenerate PPT/PDF with the new content
        return self.profile
    
