import re
import os
import io
import json
import pdfkit
import markdown
import CompanyProfile
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.utility import get_logo

class ProfileGenerator:
    def __init__(self, company_name, company_website=None):
        self.profile_data = None
        self.markdown_data = None
        self.company_name = company_name
        self.company_website = company_website
        logo_name = company_name.lower().replace(' ', '').replace('.', '')
        self.logo_url = f"https://logo.clearbit.com/{logo_name}.com"
        self.logo = get_logo(self.logo_url)
    

    def generate_markdown(self):
        markdown_text = ""
        try:
            if self.profile_data is not None:
                # Company Overview
                if 'overview' in self.profile_data:
                    overview = self.profile_data['overview']

                    company_name = overview.get('name', self.company_name)
                    logo_url = self.logo_url if self.logo is not None else overview['logo']

                    self.profile_data['overview']['logo'] = logo_url
                    
                    markdown_text += f"# <img src='{logo_url}' height='50' align='left' style='margin-right: 10px;'> {company_name} \n\n"
                    markdown_text += f"{overview.get('description', '')}\n\n"
                    
                    for key in ['mission', 'vision', 'founded', 'industry', 'location', 'employees', 'certifications']:
                        if key in overview:
                            markdown_text += f"* **{key.title()}**: {overview[key]}\n"
                    
                    markdown_text += "\n---\n\n"
                
                # Financial Highlights
                if 'financial_highlights' in self.profile_data:
                    markdown_text += "## Financial Highlights\n\n"
                    financials = self.profile_data['financial_highlights']
                    markdown_text += f"{financials.get('overview', '')}\n\n"
                    if 'metrics' in financials:
                        for year in sorted(set(financials['metrics'].keys())):
                            metric = financials['metrics'][year]
                            markdown_text += f"* **{year}**:\n"
                            markdown_text += f"  * **Revenue**: {metric['revenue']} \n"
                            markdown_text += f"  * **Growth**: {metric['growth']} \n"
                            markdown_text += f"  * **EBIT**: {metric['ebit']}\n "
                            markdown_text += f"  * **EBITDA**: {metric['ebitda']} \n"
                            markdown_text += f"  * **Gross Profit**: {metric['gross_profit']} \n"
                            markdown_text += f"  * **Net Profit**: {metric['net_profit']} \n"
                            markdown_text += f"  * **Total Assets**: {metric['total_assets']} \n"
                            markdown_text += f"  * **Market Cap**: {metric['market_cap']} \n"
                            # markdown_text += f"  * **Ownership**: {metric['ownership']} \n"
                    
                    markdown_text += "\n---\n"
                
                # Products and Services
                if 'products_services' in self.profile_data:
                    markdown_text += "## Products & Services\n\n"
                    products_services = self.profile_data['products_services']
                    markdown_text += f"{products_services.get('description', '')}\n\n"
                    for item in products_services.get('items', []):
                        markdown_text += f"#### {item['name']}\n"
                        markdown_text += f"{item['description']}\n\n"
                        markdown_text += "\n"
                    
                    markdown_text += "---\n\n"

                # Geographic Presence
                if 'geographic_presence' in self.profile_data:
                    markdown_text += "## Geographic Presence\n\n"
                    for presence in self.profile_data['geographic_presence']:
                        markdown_text += f"* **{presence['presence_type']}**: {presence['locations']}\n"
                    
                    markdown_text += "\n---\n\n"

                # Leadership
                if 'leadership' in self.profile_data:
                    markdown_text += "## Leadership Team\n\n"
                    leadership = self.profile_data['leadership']
                    markdown_text += f"{leadership.get('description', '')}\n\n"
                    for member in leadership.get('members', []):
                        markdown_text += f"#### {member['name']} - {member['position']}\n"
                        markdown_text += f"{member['bio']}\n\n"
                    
                    markdown_text += "---\n\n"

                # Strategic Priorities
                if 'strategic_priorities' in self.profile_data:
                    markdown_text += "## Strategic Priorities\n\n"
                    strategic_priorities = self.profile_data['strategic_priorities']
                    markdown_text += f"{strategic_priorities.get('description', '')}\n\n"
                    for objective in strategic_priorities.get('objectives', []):
                        markdown_text += f"* **{objective['name']}**: {objective['description']}\n"
                    
                    markdown_text += "\n---\n\n"
                
                # Clients and Competitors
                if 'clients_competitors' in self.profile_data:
                    markdown_text += "## Clients & Competitors\n\n"
                    clients_competitors = self.profile_data['clients_competitors']
                    if 'major_clients' in clients_competitors:
                        markdown_text += f"### Major Clients\n\n{clients_competitors['major_clients']}\n\n"
                    
                    if 'major_competitors' in clients_competitors:
                        markdown_text += f"### Major Competitors\n\n{clients_competitors['major_competitors']}\n\n"
                    
                    markdown_text += "\n---\n\n"

                # Key Events
                if 'key_events' in self.profile_data and self.profile_data['key_events']:
                    markdown_text += "## Key Developments\n\n"
                    key_events = self.profile_data['key_events']
                    # markdown_text += f"{key_events.get('description', '')}\n\n"
                    for event in key_events.get('events', []):
                        markdown_text += f"* **{event['date']}**: {event['description']} (Source: {event['source']})\n"
                    
                    markdown_text += "\n---\n\n"
                
                # Sources
                if 'sources' in self.profile_data:
                    markdown_text += "## Sources\n\n"
                    for source in self.profile_data['sources']:
                        markdown_text += f"* {source}\n"
                    
                    markdown_text += "\n---\n"
        except Exception as e:
            print(f"An error occurred while generating markdown: {str(e)}")
        
        return markdown_text
    

    def generate_profile(self):
        try:
            profile = CompanyProfile(self.company_name, self.company_website)
            self.profile_data = profile.get_company_profile()
            if self.profile_data is not None:
                self.markdown_data = self.generate_markdown()

            return self.markdown_data
        except Exception as e:
            print(f"An error occurred while generating company profile: {str(e)}")
        
        return None
    
    
    def generate_ppt(self):
        if self.profile_data is not None:
            try:
                prs = Presentation()
                    
                # Title slide with logo
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                title = title_slide.shapes.title
                company_name = self.profile_data['overview']['name'] if 'overview' in self.profile_data else self.company_name
                title.text = company_name

                subtitle = title_slide.placeholders[1]
                company_site = self.profile_data['overview']['website'] if 'overview' in self.profile_data else self.company_website
                subtitle.text = f"{company_site} \n\n Company Profile" if company_site else "Company Profile"
                
                logo = self.logo
                if logo is None:
                    if 'overview' in self.profile_data and 'logo' in self.profile_data['overview']:
                        try:
                            logo_url = self.profile_data['overview']['logo']
                            logo = get_logo(logo_url)
                        except Exception as e:
                            print(f"Error adding logo: {str(e)}")
                
                if logo is not None:
                    left = Inches(4)  # Center position
                    top = Inches(0.5)  # Adjusted top position
                    width = Inches(2)  # Fixed width
                    height = Inches(2)  # Fixed height
                    title_slide.shapes.add_picture(logo, left, top, width, height)

                # Company Overview slide
                if 'overview' in self.profile_data:
                    overview = self.profile_data['overview']

                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = overview['name']
                    content = slide.placeholders[1].text_frame            
                    p = content.paragraphs[0]
                    p.text = overview['description']
                    p.font.size = Pt(14)
                    
                    for key in ['industry', 'location', 'mission', 'vision']:
                        if key in overview:
                            p = content.add_paragraph()
                            p.text = f"{key.title()}: {overview[key]}"
                            p.font.size = Pt(14)
                            p.level = 0
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Financial Highlights slide
                if 'financial_highlights' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Financial Highlights"
                    
                    content = slide.placeholders[1].text_frame
                    financials = self.profile_data['financial_highlights']
                    
                    p = content.paragraphs[0]
                    p.text = financials['overview']
                    p.font.size = Pt(14)
                    
                    if 'metrics' in financials:
                        for metric in financials['metrics']:
                            p = content.add_paragraph()
                            p.text = f"{metric['year']}: Revenue {metric['revenue']}, Growth {metric['growth']}"
                            p.font.size = Pt(14)
                            p.level = 1
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Products and Services slide
                if 'products_services' in self.profile_data:
                    products_services = self.profile_data['products_services']
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Products and Services"
                    
                    content = slide.placeholders[1].text_frame
                    p = content.paragraphs[0]
                    p.text = products_services['description']
                    p.font.size = Pt(14)

                    for item in products_services['items']:
                        p = content.add_paragraph()
                        p.text = f"{item['name']}: {item['description']}"
                        p.font.size = Pt(14)
                        p.level = 0

                        if 'features' in item:
                            p = content.add_paragraph()
                            p.text = "Key Features:"
                            p.font.size = Pt(14)
                            p.level = 1
                            for feature in item['features']:
                                p = content.add_paragraph() 
                                p.text = feature
                                p.font.size = Pt(14)
                                p.level = 2
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Geographic Presence slide
                if 'geographic_presence' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Geographic Presence"
                    
                    content = slide.placeholders[1].text_frame
                    for presence in self.profile_data['geographic_presence']:
                        p = content.add_paragraph()
                        p.text = f"{presence['presence_type']}: {presence['locations']}"
                        p.font.size = Pt(14)
                        p.level = 0
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Leadership slide
                if 'leadership' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Leadership"
                    
                    content = slide.placeholders[1].text_frame
                    leadership = self.profile_data['leadership']
                    
                    p = content.paragraphs[0]
                    p.text = leadership['description']
                    p.font.size = Pt(14)
                    
                    for member in leadership['members']:
                        p = content.add_paragraph()
                        p.text = f"{member['name']} - {member['position']}"
                        p.font.size = Pt(14)
                        p.level = 0
                        
                        p = content.add_paragraph()
                        p.text = member['bio']
                        p.font.size = Pt(14)
                        p.level = 1
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Strategic Priorities slide
                if 'strategic_priorities' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Strategic Priorities"
                    
                    content = slide.placeholders[1].text_frame
                    strategic_priorities = self.profile_data['strategic_priorities']
                    
                    p = content.paragraphs[0]
                    p.text = strategic_priorities['description']
                    p.font.size = Pt(14)
                    
                    for objective in strategic_priorities['objectives']:
                        p = content.add_paragraph()
                        p.text = f"{objective['name']}: {objective['description']}"
                        p.font.size = Pt(14)
                        p.level = 1
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Clients and Competitors slide
                if 'clients_competitors' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Clients & Competitors"
                    
                    content = slide.placeholders[1].text_frame
                    clients_competitors = self.profile_data['clients_competitors']
                    
                    if 'major_clients' in clients_competitors:
                        p = content.paragraphs[0]
                        p.text = "Major Clients"
                        p.font.size = Pt(14)
                        
                        p = content.add_paragraph()
                        p.text = clients_competitors['major_clients']
                        p.font.size = Pt(14)
                        p.level = 1
                    
                    if 'major_competitors' in clients_competitors:
                        p = content.add_paragraph()
                        p.text = "Major Competitors"
                        p.font.size = Pt(14)
                        
                        p = content.add_paragraph()
                        p.text = clients_competitors['major_competitors']
                        p.font.size = Pt(14)
                        p.level = 1
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Key Events slide
                if 'key_events' in self.profile_data and self.profile_data['key_events']:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Key Events"

                    content = slide.placeholders[1].text_frame
                    key_events = self.profile_data['key_events']
                    
                    #p = content.paragraphs[0]
                    #p.text = key_events['description']
                    #p.font.size = Pt(14)
                    
                    for event in key_events['events']:
                        p = content.add_paragraph()
                        p.text = f"{event['date']}: {event['description']} (Source: {event['source']})"
                        p.font.size = Pt(14)
                        p.level = 1
                    
                    prs.slides.add_slide(prs.slide_layouts[5])  # Add a blank slide for separation

                # Sources slide
                if 'sources' in self.profile_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Sources"
                    
                    content = slide.placeholders[1].text_frame
                    for source in self.profile_data['sources']:
                        p = content.add_paragraph()
                        p.text = source
                        p.font.size = Pt(14)
                        p.level = 0

                # Save to bytes buffer
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                return ppt_buffer
            except Exception as e:
                print(f"An error occurred while generating PPT: {str(e)}")
        
        return None
    

    def generate_pdf(self):
        try:
            if self.markdown_data is not None:
                wkhtmltopdf_path = os.getenv('WKHTMLTOPDF_PATH')
                config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)

                # Convert markdown to HTML
                html = markdown.markdown(self.markdown_data)
                
                # Parse with BeautifulSoup to ensure proper HTML structure
                soup = BeautifulSoup(html, 'html.parser')
                
                # Add basic styling
                style = soup.new_tag('style')
                style.string = """
                    body { font-family: Arial, sans-serif; line-height: 1.6; padding: 20px; }
                    h1, h2 { color: #333; }
                    hr { border: 1px solid #ddd; }
                    ul { margin-left: 20px; }
                """
                soup.insert(0, style)
                
                # Convert back to string
                html_content = str(soup)
                pdf_buffer = pdfkit.from_string(html_content, configuration=config)
                return pdf_buffer
        except Exception as e:
            print(f"An error occurred while generating PDF: {str(e)}")

        return None
    

    def format_profile(self, profile_data):
        """Format an existing profile from data"""
        self.profile_data = profile_data
        self.markdown_data = self.generate_markdown()

        return self.markdown_data
    

    def update_profile_with_markdown(self, markdown_content):
        """Update the profile data based on edited markdown content"""
        # Store the original markdown
        self.profile = markdown_content
        
        # Here you could add logic to parse the markdown and update specific sections
        # of the profile_data if needed. For now, we'll just keep the markdown updated
        # and regenerate PPT/PDF with the new content
        return self.profile
    
